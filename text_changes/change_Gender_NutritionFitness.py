import os
import json
import random
import pandas as pd
from pptx.util import Cm, Pt
from pptx import Presentation
from config import MALE_FITNESS_EXCELL, FEMALE_FITNESS_EXCELL
from pptx_replace import replace_text  # Using pptx-replace

def update_gender_nutrition_fitness(json_path, ppt_path):
    """
    Updates the 36th slide of the PowerPoint with gender-based Do's and Don'ts.
    """
    # Load JSON data
    with open(json_path, 'r') as file:
        data = json.load(file)

    # Ensure case-insensitive access to keys
    data = {k.lower(): v for k, v in data.items()}

    age = data.get("age")  # Safely access the 'age' key
    gender = data.get("gender")

    if age is None or gender is None:
        raise ValueError("Missing 'age' or 'gender' in the JSON data.")

    # Select the appropriate Excel file
    excel_file = FEMALE_FITNESS_EXCELL if gender.lower() == "women" else MALE_FITNESS_EXCELL

    # Load the Excel file and normalize column headers
    df = pd.read_excel(excel_file)
    df.columns = [col.replace("’", "'") for col in df.columns]  # Replace curly apostrophes with straight ones

    # Find the appropriate age range and extract dos and don'ts
    dos = []
    donts = []

    for _, row in df.iterrows():
        age_range = row["Age"]  # Assuming the column header is "Age"
        if "-" in str(age_range):
            # Remove non-numeric characters except for the dash
            cleaned_age_range = ''.join(c for c in age_range if c.isdigit() or c == '-')
            min_age, max_age = map(int, cleaned_age_range.split("-"))
            if min_age <= age <= max_age:
                # Append values instead of overwriting
                dos.extend(row["Do's"].split(","))  # Assuming the column header is "Dos"
                donts.extend(row["Don't's"].split(",") if pd.notna(row["Don't's"]) else [])

    # Remove duplicates and strip whitespace
    dos = list(set(d.strip() for d in dos))
    donts = list(set(d.strip() for d in donts))

    # Select random 4 points for Do's and 5 points for Don'ts
    dos = random.sample(dos, min(4, len(dos)))
    donts = random.sample(donts, min(5, len(donts)))

    # Load the PowerPoint presentation
    prs = Presentation(ppt_path)
    slide = prs.slides[35]  # 36th slide (0-based index)

    # Replace Do's placeholders (4 points)
    for i in range(1, 5):  # 1 to 4
        placeholder = f"{{dopoint{i}}}"
        text_to_insert = dos[i-1] if i-1 < len(dos) else ""
        replace_text(prs, placeholder, text_to_insert)

    # Replace Don'ts placeholders (5 points)
    for i in range(1, 6):  # 1 to 5
        placeholder = f"{{dontpoint{i}}}"
        text_to_insert = donts[i-1] if i-1 < len(donts) else ""
        replace_text(prs, placeholder, text_to_insert)

    # Save the updated PowerPoint
    prs.save(ppt_path)
