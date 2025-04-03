import os, json
from pptx import Presentation
from pptx.util import Cm, Pt
import openpyxl
from config import GENERATED_OUTPUTS, patients_folder as PF, RDA_FILE
import pandas as pd

# Define text box parameters for each Risk level
TEXT_BOX_PARAMS_RISK = {
    3: [
        (Cm(4.61), Cm(5.43), Cm(5.25), Cm(7.11)),
        (Cm(4.61), Cm(5.43), Cm(5.25), Cm(13.12)),
    ],
    2: [
        (Cm(4.61), Cm(4.58), Cm(13.07), Cm(8.7)),
        (Cm(4.61), Cm(4.58), Cm(13.07), Cm(14.2)),
    ],
    1: [
        (Cm(4.61), Cm(4.43), Cm(20.79), Cm(6.98)),
        (Cm(4.61), Cm(4.43), Cm(20.79), Cm(13.11)),
    ]
}

def add_text_boxes_on_slide(prs, slide_index, patient_code):
    VITAMIN_SHEET_FILE = os.path.join(PF, patient_code, f"{patient_code}_vitamin_sheet.xlsx")
    JSON_PATH = os.path.join(PF, patient_code, f"{patient_code}.json")

    if not os.path.exists(VITAMIN_SHEET_FILE) or not os.path.exists(RDA_FILE) or not os.path.exists(JSON_PATH):
        print(f"❌ Required files not found for patient {patient_code}")
        return False

    # Read the Excel files and JSON
    df = pd.read_excel(VITAMIN_SHEET_FILE)
    rda_df = pd.read_excel(RDA_FILE)
    
    # Get patient gender from JSON
    with open(JSON_PATH, 'r') as file:
        json_data = json.load(file)
    json_data = {k.lower(): v for k, v in json_data.items()}
    gender = json_data.get("gender", "Female")  # Default to Female if not found
    rda_column = 'Female (mg/day)' if gender.lower() == 'female' else 'Male (mg/day)'

    # Initialize dictionaries for Risk levels
    risk_dict = {3: [], 2: [], 1: []}
    risk_columns_dict = {3: set(), 2: set(), 1: set()}

    # Process the data
    for _, row in df.iterrows():
        risk_level = row['Risk']
        if risk_level in risk_dict:
            condition = row['Condition']
            # Find matching nutrient in RDA file
            rda_match = rda_df[rda_df['Nutrient'].str.contains(fr'\b{condition}\b', case=False, na=False, regex=True)]
            if not rda_match.empty:
                rda_value = rda_match.iloc[0][rda_column]
                condition_with_rda = f"{condition} ({rda_value})"
                risk_dict[risk_level].append(condition_with_rda)
            else:
                risk_dict[risk_level].append(condition)

    # Convert sets to sorted lists for consistent ordering
    risk_columns_dict = {k: sorted(v) for k, v in risk_columns_dict.items()}

    # Get the specified slide
    slide = prs.slides[slide_index]

    # Function to add text to a text box with bullet points
    def add_bullet_points(text_box, items, font_name, font_size, bold=False):
        text_frame = text_box.text_frame
        text_frame.clear()  # Clear any existing text

        for idx, item in enumerate(items):
            if idx == 0:
                # Use the default paragraph for the first bullet point
                p = text_frame.paragraphs[0]
            else:
                # Add a new paragraph for subsequent bullet points
                p = text_frame.add_paragraph()
            p.text = item
            p.font.name = font_name
            p.font.size = font_size
            p.font.bold = bold
            p.space_after = Cm(0.1)  # Adjust spacing between bullet points

    # Function to add wrapped text to a text box with specific line formatting
    def add_wrapped_text_with_lines(text_box, items, items_per_line, font_name, font_size, bold=False, italic=False):
        text_frame = text_box.text_frame
        text_frame.clear()
        lines = []
        current_line = []
        for i, item in enumerate(items, 1):
            current_line.append(item)
            if i % items_per_line == 0 or i == len(items):
                lines.append(", ".join(current_line))
                current_line = []
        text = "\n".join(lines)
        p = text_frame.paragraphs[0]
        run = p.add_run()
        run.text = text
        run.font.name = font_name
        run.font.size = font_size
        run.font.bold = bold
        run.font.italic = italic

    # Add text boxes for each Risk level
    for risk, params in TEXT_BOX_PARAMS_RISK.items():
        conditions = risk_dict[risk]
        risk_columns = risk_columns_dict[risk]

        # Split conditions into two lists for the first and second text boxes
        conditions_text_1 = []
        conditions_text_2 = []
        current_height = 0
        max_height = Cm(4)  # Approximate height limit for the first text box

        for condition in conditions:
            if current_height + Cm(0.5) <= max_height:  # Approximate height per bullet point
                conditions_text_1.append(condition)
                current_height += Cm(0.5)
            else:
                conditions_text_2.append(condition)

        # Determine items per line for the third box based on risk level
        items_per_line = 3 if risk == 3 else 2

        for i, (height, width, top, left) in enumerate(params):
            text_box = slide.shapes.add_textbox(left, top, width, height)
            if i == 0:  # First box
                add_bullet_points(
                    text_box,
                    conditions_text_1,
                    font_name="Arial",
                    font_size=Pt(11),
                    bold=True
                )
            elif i == 1:  # Second box
                add_bullet_points(
                    text_box,
                    conditions_text_2,
                    font_name="Arial",
                    font_size=Pt(11),
                    bold=True
                )

def update_vitamin_details(patient_id):
    ppt_path = os.path.join(GENERATED_OUTPUTS, f"{patient_id}_report.pptx")

    prs = Presentation(ppt_path)
    if add_text_boxes_on_slide(prs, slide_index=38, patient_code=patient_id):  # Add to the 39th slide (index 38)
        prs.save(ppt_path)
        print(f"✅ Vitamin details updated for patient {patient_id}")

    prs = Presentation(ppt_path)
    add_text_boxes_on_slide(prs, slide_index=38, patient_code=patient_id)  # Add to the 39th slide (index 38)
    prs.save(ppt_path)
    print(f"✅ Vitamin details updated for patient {patient_id}")
