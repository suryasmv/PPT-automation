import os
import pandas as pd
from pptx import Presentation
from pptx.util import Cm,Pt
from config import scoring_charts, DIET_PICTURES as IMAGE_PATH, DIET_FILE

START_SLIDE = 26  # 0-based index (27th slide)
END_SLIDE = 29    # 0-based index (30th slide)
OTHER_START_SLIDE = 30  # 0-based index (31st slide)
OTHER_END_SLIDE = 34    # 0-based index (35th slide)
START_X = Cm(1)
START_Y = Cm(9)
END_Y = Cm(27)
CARD_GAP = Cm(1)
SEVERITY_ORDER = ["Moderate to High", "Moderate", "Mild"]  # Priority order

# Specific Moderate Conditions
MODERATE_SPECIFIC_CONDITIONS = {
    "Cardiac_Health",
    "Cholesterol_Disorders",
    "High_Blood_Pressure",
    "Glomerular_Diseases",
    "Allergies",
    "Gut_Health",
}

MODERATE_TO_HIGH_SPECIFIC_CONDITIONS = {
    "Allergies",
    "Cardiac_Health",
    "Cardiomyopathy",
    "Cholesterol_Disorders",
    "Diabetes",
    "Gut_Health",
    "High_Blood_Pressure",
    "Obesity",
    "Stroke"
}

def find_scoring_chart(patient_id):
    for root, _, files in os.walk(scoring_charts):
        for file in files:
            if file.startswith(f"{patient_id}_Scoring_chart") and file.endswith(".xlsx"):
                return os.path.join(root, file)
    return None

def extract_severity_conditions(excel_path):
    df = pd.read_excel(excel_path)
    concern_conditions = []
    other_conditions = []
    
    for _, row in df.iterrows():
        condition_name = row["Medical Condition "].replace(" ", "_")
        is_concern = str(row.get("concerns", "")).strip().lower() == 'y'
        
        for severity in SEVERITY_ORDER:
            if severity in df.columns and str(row.get(severity, '')).strip().lower() == 'y':
                if is_concern:
                    concern_conditions.append((severity, condition_name))
                else:
                    other_conditions.append((severity, condition_name))
                    
    return concern_conditions, other_conditions

def extract_recommendations(condition, severity):
    """
    Extracts recommendation points for a given condition and severity from the DIET_FILE.
    """
    df = pd.read_excel(DIET_FILE)
    if severity in df.columns:
        recommendations = df[df["Condition"] == condition][severity].dropna().tolist()
        formatted_recommendations = []
        for rec in recommendations:
            points = rec.split('$')
            for point in points:
                if point.strip():
                    formatted_recommendations.append(f"\u2022 {point.strip().capitalize()}")
        return "\n".join(formatted_recommendations)
    return ""

def add_recommendation_textbox(slide, x_pos, y_pos, card_width, card_height, condition, severity):
    """
    Adds a recommendation text box to the slide based on the card type and dimensions.
    """
    if severity == "Mild":
        tb_height, tb_width = Cm(3.8), Cm(16.09)
        tb_x, tb_y = x_pos + Cm(2.6), y_pos + Cm(1)
    elif severity == "Moderate" and condition not in MODERATE_SPECIFIC_CONDITIONS:
        tb_height, tb_width = Cm(3.7), Cm(15.98)
        tb_x, tb_y = x_pos + Cm(2.7), y_pos + Cm(1)
    elif severity == "Moderate to High" and condition not in MODERATE_TO_HIGH_SPECIFIC_CONDITIONS:
        tb_height, tb_width = Cm(4), Cm(16.25)
        tb_x, tb_y = x_pos + Cm(2.5), y_pos + Cm(1.15)
    else:
        if condition in MODERATE_SPECIFIC_CONDITIONS:
            tb_height, tb_width = Cm(5.3), Cm(15.98)
            tb_x, tb_y = x_pos + Cm(2.5), y_pos + Cm(1)
            severity = "Moderate"
        if condition in MODERATE_TO_HIGH_SPECIFIC_CONDITIONS:
            tb_height, tb_width = Cm(6.5), Cm(16.25)
            tb_x, tb_y = x_pos + Cm(2.5), y_pos + Cm(1.15)  
            severity == "Moderate to High"      

    text_box = slide.shapes.add_textbox(tb_x, tb_y, tb_width, tb_height)
    text_frame = text_box.text_frame
    text_frame.clear()  # Clear default placeholder text

    recommendations = extract_recommendations(condition, severity)
    text_frame.word_wrap = True
    for idx, rec in enumerate(recommendations.split('\n')):
        if idx == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        p.text = rec
        p.font.name = "Arial"
        # Adjust font size based on text length
        p.font.size = Pt(11)

def add_diet_images(ppt_file, patient_id):
    """
    Adds diet images to the specified slides in the PowerPoint file.
    """
    # Find scoring chart
    scoring_chart = find_scoring_chart(patient_id)
    if not scoring_chart:
        print(f"Scoring chart not found for patient {patient_id}.")
        return

    # Extract severity conditions
    concern_conditions, other_conditions = extract_severity_conditions(scoring_chart)
    if not concern_conditions and not other_conditions:
        print(f"No conditions found for patient {patient_id}.")
        return

    # Load presentation
    prs = Presentation(ppt_file)

    # Validate slide count
    if len(prs.slides) <= OTHER_END_SLIDE:
        print("PPTX file doesn't have enough slides.")
        return

    def find_image(condition, severity_folder):
        """Finds the image path based on the condition and severity."""
        image_folder = os.path.join(IMAGE_PATH, severity_folder)
        for file in os.listdir(image_folder):
            if condition.lower() in file.lower():
                return os.path.join(image_folder, file)
        return None

    # Updated Dimensions
    MODERATE_SPECIFIC_HEIGHT = Cm(6.84)
    MODERATE_SPECIFIC_WIDTH = Cm(19.17)

    OTHER_MODERATE_HEIGHT = Cm(5.04)
    OTHER_MODERATE_WIDTH = Cm(19.18)

    MODERATE_TO_HIGH_HEIGHT = Cm(5.54)
    MODERATE_TO_HIGH_WIDTH = Cm(19.18)

    OTHER_MODERATE_TO_HIGH_HEIGHT = Cm(8.26)
    OTHER_MODERATE_TO_HIGH_WIDTH = Cm(19.13)

    MILD_HEIGHT = Cm(5.25)
    MILD_WIDTH = Cm(19.16)

    def insert_images(conditions, start_slide, end_slide):
        slide_index = start_slide
        x_pos, y_pos = START_X, START_Y

        # Sort conditions by severity order
        conditions.sort(key=lambda x: SEVERITY_ORDER.index(x[0]) if x[0] in SEVERITY_ORDER else float('inf'))

        for severity, condition in conditions:
            if severity == "Mild":
                card_height, card_width = MILD_HEIGHT, MILD_WIDTH
            elif severity == "Moderate":
                if condition in MODERATE_SPECIFIC_CONDITIONS:
                    card_height, card_width = MODERATE_SPECIFIC_HEIGHT, MODERATE_SPECIFIC_WIDTH
                else:
                    card_height, card_width = OTHER_MODERATE_HEIGHT, OTHER_MODERATE_WIDTH
            elif severity == "Moderate to High":
                if condition in MODERATE_SPECIFIC_CONDITIONS:
                    card_height, card_width = OTHER_MODERATE_TO_HIGH_HEIGHT, OTHER_MODERATE_TO_HIGH_WIDTH
                else:
                    card_height, card_width = MODERATE_TO_HIGH_HEIGHT, MODERATE_TO_HIGH_WIDTH
            else:
                # Treat unhandled severity levels as "Moderate" for specific conditions
                if condition in MODERATE_SPECIFIC_CONDITIONS:
                    card_height, card_width = MODERATE_SPECIFIC_HEIGHT, MODERATE_SPECIFIC_WIDTH
                elif condition in MODERATE_TO_HIGH_SPECIFIC_CONDITIONS:
                    card_height, card_width = OTHER_MODERATE_TO_HIGH_HEIGHT, OTHER_MODERATE_TO_HIGH_WIDTH
                else:
                    print(f"Unhandled severity level: {severity} for condition: {condition}")
                    continue  # Skip if not in specific conditions

            severity_folder = "Mild" if severity == "Mild" else "Moderate" if severity == "Moderate" else "Moderate_to_High"
            
            image_path = find_image(condition, severity_folder)
            if not image_path:
                print(f"Image not found for {condition} in {severity_folder} folder.")
                continue
            
            slide = prs.slides[slide_index]
            slide.shapes.add_picture(image_path, x_pos, y_pos, width=card_width, height=card_height)

            # Add recommendation text box
            add_recommendation_textbox(slide, x_pos, y_pos, card_width, card_height, condition, severity)
            
            # Update Y position
            y_pos += card_height + CARD_GAP
            
            # Check if we exceed the slide limit
            if y_pos + card_height > END_Y:
                slide_index += 1
                y_pos = START_Y
                if slide_index > end_slide:
                    print("No more space in slides.")
                    break

    # Insert concern conditions
    insert_images(concern_conditions, START_SLIDE, END_SLIDE)
    # Insert other conditions
    insert_images(other_conditions, OTHER_START_SLIDE, OTHER_END_SLIDE)

    prs.save(ppt_file)
    print("Diet image insertion completed.")
