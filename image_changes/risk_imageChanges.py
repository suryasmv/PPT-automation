import os
import shutil
import pandas as pd
from pptx import Presentation
from pptx.util import Inches
from config import scoring_charts, image_paths

def process_excel(patient_code):
    """Extract severity levels for medical conditions from the patient's Excel file."""
    excel_path = os.path.join(scoring_charts, f"BATCH1/{patient_code}_Scoring_chart.xlsx")
    df = pd.read_excel(excel_path)
    
    severity_columns = ["Low", "Mild", "Moderate", "Moderate to High"]
    severity_mapping = {}
    
    for _, row in df.iterrows():
        condition = row["Medical Condition "]
        severity = next((col for col in severity_columns if row.get(col) == 'y'), "Low")
        severity_mapping[condition] = severity
    
    print("\n✅ Extracted Severity Mapping:", severity_mapping)
    return severity_mapping

def replace_ppt_images(ppt_path, severity_mapping):
    """Replace images in PowerPoint slides based on severity levels."""
    prs = Presentation(ppt_path)
    slide_mapping = {
        6: [
            "Diabetes", "High_Blood_Pressure", "Cardiac_Health", "Cholesterol_Disorders",
            "Cardiomyopathy", "Arrhythmias", "Obesity", "Thyroid_Disorders",
            "Dementia", "Stroke", "Glomerular_Diseases"
        ],
        7: [
            "Mood_Disorders", "Fatty_Liver", "Gall_stones", "Gastritis",
            "Gut_Health", "Allergies", "Skin_Health", "Muscular_health"
        ]
    }
    
    height_inches = 1.36 * 0.3937  # Convert cm to inches
    width_inches = 1.46 * 0.3937
    
    for slide_index, conditions in slide_mapping.items():
        slide = prs.slides[slide_index]
        image_shapes = [shape for shape in slide.shapes if shape.shape_type == 13]
        
        for i, condition in enumerate(conditions):
            if i < len(image_shapes):
                severity = severity_mapping.get(condition, "Low")
                new_image_path = image_paths[severity]
                shape = image_shapes[i]
                
                slide.shapes._spTree.remove(shape._element)
                slide.shapes.add_picture(new_image_path, shape.left, shape.top, width=Inches(width_inches), height=Inches(height_inches))
                print(f"✅ Replaced image for '{condition}' on slide {slide_index + 1} with severity '{severity}'.")
    
    prs.save(ppt_path)
    print(f"\n✅ Updated PowerPoint saved: {ppt_path}")

def process_risk_images(patient_code, ppt_path):
    """Main function to process risk images for a given patient."""
    severity_mapping = process_excel(patient_code)
    replace_ppt_images(ppt_path, severity_mapping)
