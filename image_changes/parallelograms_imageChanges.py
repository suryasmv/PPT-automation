import os
import pandas as pd
from pptx import Presentation
from pptx.util import Cm,Pt
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from config import scoring_charts, input_parallelograms, generated_outputs as GO, RECOMMENDATIONS_FILE, FIRST_TEXT_FILE

# Image sizes
IMAGE_SIZES = {
    "Moderate to High": (Cm(9.35), Cm(19.43)),
    "Moderate": (Cm(7), Cm(19.43)),
    "Mild": (Cm(7), Cm(19.43)),
    "Low": (Cm(5), Cm(19.43))
}

LOW_LIST = [
    # "Anemia",
    "Arrhythmias",
    "Cardiac_Health",
    "Cardiomyopathy",
    "Cholesterol_Disorders",
    "Dementia",
    "Diabetes",
    "Fatty_Liver",
    "Gall_stones",
    "Gastritis",
    "Glomerular_Diseases",
    "Gut_Health",
    # "Headaches",
    "High_Blood_Pressure",
    "Mood_Disorders",
    "Muscular_health",
    "Obesity",
    # "Pancreatic_Disorders",
    # "Parkinsons",
    "Renal_stones",
    "Allergies",
    "Skin_Health",
    "Stroke",
    "Thyroid_Disorders"
]

# Text box sizes and positions
TEXT_BOX_PARAMS = {
    "Moderate": [
        (Cm(1.38), Cm(15.36), Cm(2.8), Cm(1.23)),  # 1st text box
        (Cm(2.71), Cm(10), Cm(2.98), Cm(3.25))  # 2nd text box (Recommendations)
    ],
    "Mild": [
        (Cm(1.32), Cm(15.38), Cm(2.8), Cm(1.27)),  # 1st text box
        (Cm(2.56), Cm(9.89), Cm(2.8), Cm(3.53))  # 2nd text box (Recommendations)
    ],
    "Moderate to High": [
        (Cm(1.32), Cm(15.38), Cm(2.8), Cm(1.60)),  # 1st text box
        (Cm(4.1), Cm(9.96), Cm(2.8), Cm(3.80))  # 2nd text box (Recommendations)
    ]
}

# Slide parameters
START_X = Cm(0.7)
START_Y = Cm(4.5)
MAX_Y = Cm(27)
SPACING = Cm(0.5)
CONCERN_START_SLIDE = 9
CONCERN_END_SLIDE = 14
OTHER_START_SLIDE = 15
OTHER_END_SLIDE = 22
SEVERITY_ORDER = ["Moderate to High", "Moderate", "Mild", "Low"]  # Priority order

BOLD_WORDS = ['Moderate', 'Mild', 'Moderate to High']

def find_scoring_chart(patient_id):
    for root, _, files in os.walk(scoring_charts):
        for file in files:
            if file.startswith(f"{patient_id}_Scoring_chart") and file.endswith(".xlsx"):
                return os.path.join(root, file)
    return None


def extract_severity_conditions(excel_path):
    df = pd.read_excel(excel_path)
    concern_conditions = []
    other_conditions = []

    for _, row in df.iterrows():
        condition_name = row["Medical Condition "].replace(" ", "_")
        is_concern = str(row.get("concerns", "")).strip().lower() == 'y'

        for severity in SEVERITY_ORDER:
            if severity in df.columns and str(row.get(severity, '')).strip().lower() == 'y':
                if is_concern:
                    concern_conditions.append((severity, condition_name))
                else:
                    other_conditions.append((severity, condition_name))
    
    print("Concern Conditions:", concern_conditions)
    print("Other Conditions:", other_conditions)
    return concern_conditions, other_conditions


def find_condition_image(severity, condition):
    severity_path = os.path.join(input_parallelograms, severity)
    if os.path.exists(severity_path):
        for file in os.listdir(severity_path):
            if file.lower().startswith(condition.lower()) and file.endswith(".png"):
                return os.path.join(severity_path, file)
    return None


def extract_recommendations(condition, severity):
    df = pd.read_excel(RECOMMENDATIONS_FILE)
    if severity in df.columns:
        recommendations = df[df["Condition"] == condition][severity].dropna().tolist()
        formatted_recommendations = []
        for rec in recommendations:
            points = rec.split('$')
            for point in points:
                if point.strip():
                    formatted_recommendations.append(f"\u2022 {point.strip()}")
        return "\n".join(formatted_recommendations)
    return ""

def extract_first_text(condition, severity):
    df = pd.read_excel(FIRST_TEXT_FILE)
    if severity in df.columns:
        first_text = df[df["Condition"] == condition][severity].dropna().tolist()
        if first_text:
            return first_text[0]
    return ""

def add_text_with_formatting(text_frame, text):
    p = text_frame.paragraphs[0]
    p.clear()
    run = p.add_run()
    run.font.name = "Arial"
    run.font.size = Pt(9)  # Arial 9

    words = text.split()
    for word in words:
        run = p.add_run()
        run.text = word + ' '
        run.font.name = "Arial"
        run.font.size = Pt(9)
        if word.strip(',') in BOLD_WORDS:
            run.font.bold = True


def insert_parallelogram_images(patient_id):
    output_ppt_path = os.path.join(GO, f"{patient_id}_report.pptx")
    if not os.path.exists(output_ppt_path):
        print(f"❌ Report not found for patient {patient_id}")
        return

    prs = Presentation(output_ppt_path)
    scoring_chart = find_scoring_chart(patient_id)
    if not scoring_chart:
        print(f"❌ No scoring chart found for patient {patient_id}")
        return

    concern_conditions, other_conditions = extract_severity_conditions(scoring_chart)
    if not concern_conditions and not other_conditions:
        print(f"⚠️ No conditions found for patient {patient_id}")
        return

    # Track processed conditions
    processed_conditions = set()

    def insert_conditions(conditions, start_slide, end_slide, start_y):
        current_y = start_y
        slide_index = start_slide
        slide = prs.slides[slide_index]

        # Sort conditions by severity order
        conditions.sort(key=lambda x: SEVERITY_ORDER.index(x[0]))

        for severity, condition in conditions:
            processed_conditions.add(condition)

            image_path = find_condition_image(severity, condition)
            if not image_path:
                print(f"❌ Image not found for {condition} ({severity})")
                continue

            img_height, img_width = IMAGE_SIZES[severity]
            if current_y + img_height > MAX_Y:
                slide_index += 1
                if slide_index > end_slide:
                    print(f"⚠️ Slide limit reached for current section (limit: {end_slide}), stopping.")
                    break
                slide = prs.slides[slide_index]
                current_y = START_Y

            slide.shapes.add_picture(image_path, START_X, current_y, width=img_width, height=img_height)

            if severity in TEXT_BOX_PARAMS:
                for i, (tb_h, tb_w, hp, vp) in enumerate(TEXT_BOX_PARAMS[severity]):
                    text_box = slide.shapes.add_textbox(START_X + hp, current_y + vp, tb_w, tb_h)
                    text_frame = text_box.text_frame
                    text_frame.clear()  # Clear default placeholder text

                    if i == 0:  # First text box
                        first_text = extract_first_text(condition, severity)
                        add_text_with_formatting(text_frame, first_text)
                    elif i == 1:  # Recommendation text box
                        text_frame.word_wrap = True
                        text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
                        recommendations = extract_recommendations(condition, severity)
                        for idx, rec in enumerate(recommendations.split('\n')):
                            if idx == 0:
                                p = text_frame.paragraphs[0]
                            else:
                                p = text_frame.add_paragraph()
                            p.text = rec
                            p.font.name = "Arial"
                            p.font.size = Pt(9)

            current_y += img_height + SPACING

        return slide_index, current_y

    # Insert concern conditions in slides 9-14
    _, _ = insert_conditions(concern_conditions, CONCERN_START_SLIDE, CONCERN_END_SLIDE, START_Y)
    
    # Insert other conditions in slides 15-29
    current_slide, current_y = insert_conditions(other_conditions, OTHER_START_SLIDE, OTHER_END_SLIDE, START_Y)

    # Add missing low conditions to the remaining space in other conditions section
    missing_low_conditions = [(severity, condition) for severity, condition in [("Low", c) for c in LOW_LIST] 
                            if condition not in processed_conditions]
    insert_conditions(missing_low_conditions, current_slide, OTHER_END_SLIDE, current_y)

    prs.save(output_ppt_path)
    print(f"✅ Parallelogram Images and Textboxes inserted for patient {patient_id}")



