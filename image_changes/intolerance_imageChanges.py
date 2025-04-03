import json
import os
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from config import patients_folder, image_paths, GENERATED_OUTPUTS

def cm_to_emu(cm):
    return int(cm * 360000)

def add_intolerance_details(patient_code):
    # Read intolerance data
    json_path = os.path.join(patients_folder, patient_code, f"{patient_code}_intolerance.json")
    with open(json_path, 'r') as f:
        data = json.load(f)
    
    ppt_path = os.path.join(GENERATED_OUTPUTS, f"{patient_code}_report.pptx")
    prs = Presentation(ppt_path)
    
    # Configuration for positions
    image_config = {
        37: [(4.62, "Carbohydrate_Intolerance"),
             (12.4, "Lipid_Intolerance"),
             (21.09, "Protein_Intolerance")],
        38: [(5.29, "Lactose_Intolerance"),
             (12.88, "Gluten_Intolerance"),
             (20.53, "Insulin_Resistance")]
    }
    
    text_config = {
        37: [(8.91, "Carbohydrate_Intolerance"),
             (16.63, "Lipid_Intolerance"),
             (25.38, "Protein_Intolerance")],
        38: [(9.58, "Lactose_Intolerance"),
             (17.17, "Gluten_Intolerance"),
             (24.82, "Insulin_Resistance")]
    }
    
    # Standard dimensions
    img_width, img_height = 1.54, 5.59
    txt_width, txt_height = 2.41, 0.77
    img_left = 16
    txt_left = 17.2
    
    for slide_num in [37, 38]:
        slide = prs.slides[slide_num - 1]
        
        # Add images
        for top, key in image_config[slide_num]:
            value = data[key]
            img_path = image_paths[f"{value}1"]
            slide.shapes.add_picture(
                img_path,
                Cm(img_left),
                Cm(top),
                width=Cm(img_width),
                height=Cm(img_height)
            )
        
        # Add text boxes
        for top, key in text_config[slide_num]:
            value = data[key]
            txBox = slide.shapes.add_textbox(
                Cm(txt_left),
                Cm(top),
                Cm(txt_width),
                Cm(txt_height)
            )
            paragraph = txBox.text_frame.paragraphs[0]
            run = paragraph.add_run()
            run.text = value
            run.font.name = 'Arial'
            run.font.size = Pt(12)
            run.font.bold = True
            run.font.color.rgb = RGBColor(255, 0, 0)
    
    prs.save(ppt_path)
