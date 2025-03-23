import os
import shutil
from pptx import Presentation
from config import patients_folder as PF, generated_outputs as GO, lifeStyle_template as LT
from text_changes.change_SequencingDetails import replace_text_in_ppt
from image_changes.parallelograms_imageChanges import insert_parallelogram_images
from image_changes.diet_imageChanges import add_diet_images  
from text_changes.change_VitaminDetails import update_vitamin_details
from image_changes.risk_imageChanges import process_risk_images
from text_changes.change_Gender_NutritionFitness import update_gender_nutrition_fitness

# Define patient names to process (one by one)
selected_patients = ["2400182MGI"]

def copy_template_ppt(target_path):
    """Copies the actual PPT template to the target location for modifications."""
    shutil.copy(LT, target_path)  
    print(f"📂 Template copied to: {target_path}")

# Define Y-coordinate bounds for checking empty slides (1 cm = 360000 EMUs)
START_Y = 9 * 360000  
MAX_Y = 26 * 360000  

def has_content(slide):
    """Checks if a slide contains any shapes (images or textboxes) in the valid area."""
    for shape in slide.shapes:
        if hasattr(shape, "top") and START_Y <= shape.top <= MAX_Y:
            return True
    return False

def delete_empty_slides(output_ppt_path):
    """Deletes all empty slides in the PowerPoint if they have no images or textboxes."""
    prs = Presentation(output_ppt_path)
    empty_slide_indexes = [i for i in range(len(prs.slides)) if not has_content(prs.slides[i])]

    if empty_slide_indexes:
        for i in sorted(empty_slide_indexes, reverse=True):
            xml_slides = prs.slides._sldIdLst
            xml_slides.remove(xml_slides[i])
        prs.save(output_ppt_path)
        print(f"🗑️ Empty slides removed: {empty_slide_indexes}")
    else:
        print(f"✅ No empty slides found.")
    
    return empty_slide_indexes 

def generate_patient_report(patient_code):
    """
    Generates a PowerPoint report for a single patient.
    """
    json_path = os.path.join(PF, patient_code, f"{patient_code}.json")
    output_ppt_path = os.path.join(GO, f"{patient_code}_report.pptx")

    if not os.path.exists(json_path):
        print(f"⚠️ JSON file not found for patient {patient_code}, skipping.")
        return
    
    print(f"\n🚀 Generating report for: {patient_code}")

    # Step 1: Copy the actual template PPT to the output folder
    copy_template_ppt(output_ppt_path)

    # Step 2: Process text replacement
    replace_text_in_ppt(json_path, output_ppt_path, output_ppt_path)
    print("📜 Text replacements done.")

    # Step 3: Insert parallelogram images
    insert_parallelogram_images(patient_code)
    print("🖼️ Parallelogram images added.")

    # Step 4: Insert diet images
    add_diet_images(output_ppt_path, patient_code)
    print("🥗 Diet images added.")

    # Step 5: Update vitamin details
    update_vitamin_details(patient_code)
    print("💊 Vitamin details updated.")

    # Step 6: Process risk images
    process_risk_images(patient_code, output_ppt_path)
    print("📊 Risk images processed.")

    # Step 7: Update gender-based nutrition and fitness details
    update_gender_nutrition_fitness(json_path, output_ppt_path)
    print("🏋️ Gender-based nutrition & fitness details updated.")

    # Step 8: Delete empty slides
    delete_empty_slides(output_ppt_path)

    print(f"✅ Report completed: {output_ppt_path}\n")


def generate_reports():
    """
    Runs the report generation process sequentially for each selected patient.
    """
    for patient in selected_patients:
        generate_patient_report(patient)

# Execute the entire process
generate_reports()
